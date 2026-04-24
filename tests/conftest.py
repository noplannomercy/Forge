import io
import pytest
from docx import Document as DocxDocument
from openpyxl import Workbook
from pptx import Presentation

from job_store import InMemoryJobStore


@pytest.fixture
def store():
    """테스트용 InMemoryJobStore — test_worker, test_callback_field_map 등에서 공유."""
    return InMemoryJobStore()


@pytest.fixture
def sample_docx_bytes():
    """텍스트 + 표가 있는 간단한 DOCX"""
    doc = DocxDocument()
    doc.add_heading("제목", level=1)
    doc.add_paragraph("본문 텍스트입니다.")
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "이름"
    table.cell(0, 1).text = "나이"
    table.cell(1, 0).text = "홍길동"
    table.cell(1, 1).text = "30"
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


@pytest.fixture
def empty_docx_bytes():
    """빈 DOCX"""
    doc = DocxDocument()
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


@pytest.fixture
def sample_pptx_bytes():
    """슬라이드 2개짜리 PPTX"""
    prs = Presentation()
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "슬라이드 1 제목"
    slide1.placeholders[1].text = "슬라이드 1 본문"
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "슬라이드 2 제목"
    slide2.placeholders[1].text = "슬라이드 2 본문"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


@pytest.fixture
def sample_xlsx_bytes():
    """시트 2개짜리 XLSX"""
    wb = Workbook()
    ws1 = wb.active
    ws1.title = "매출"
    ws1.append(["월", "금액"])
    ws1.append(["1월", 1000])
    ws1.append(["2월", 2000])
    ws2 = wb.create_sheet("비용")
    ws2.append(["항목", "금액"])
    ws2.append(["인건비", 500])
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()
