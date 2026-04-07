import io
from pptx import Presentation

from models import ConvertResult, Quality


async def extract(file_bytes: bytes, file_name: str) -> ConvertResult:
    """PPTX → Markdown 변환 (슬라이드별)"""
    prs = Presentation(io.BytesIO(file_bytes))
    parts: list[str] = []

    for i, slide in enumerate(prs.slides, 1):
        slide_parts: list[str] = [f"## 슬라이드 {i}"]

        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_parts.append(text)
            if shape.has_table:
                table = shape.table
                rows = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    rows.append("| " + " | ".join(cells) + " |")
                if rows:
                    header_sep = "| " + " | ".join(["---"] * len(table.columns)) + " |"
                    rows.insert(1, header_sep)
                    slide_parts.append("\n".join(rows))

        parts.append("\n\n".join(slide_parts))

    full_text = "\n\n---\n\n".join(parts)
    total_chars = len(full_text.strip())
    num_slides = len(prs.slides)

    return ConvertResult(
        text=full_text,
        format="md",
        pages=num_slides,
        file_name=file_name,
        source_format="pptx",
        route="extract",
        quality=Quality(
            total_chars=total_chars,
            chars_per_page=total_chars / num_slides if num_slides > 0 else 0,
            total_pages=num_slides,
            failed_pages=0,
            confidence="high",
        ),
    )
